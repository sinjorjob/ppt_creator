from pptx import Presentation 
from pptx.util import Inches 
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Cm
from pptx.util import Pt
from pptx.util import Inches
from pptx.dml.color import RGBColor
def addImage(slide, url):
    left = top = Inches(1)
    slide.shapes.add_picture(url, left, top)

""" Ref for slide types: 
0 ->  title and subtitle
1 ->  title and content
2 ->  section header
3 ->  two content
4 ->  Comparison
5 ->  Title only 
6 ->  Blank
7 ->  Content with caption
8 ->  Pic with caption
"""

def add_slide(presentation, type):
    layout = presentation.slide_layouts[type]  #スライドレイアウトの選択
    return presentation.slides.add_slide(layout) #スライドの追加


def add_splash_slide(presentation, title, subtitle):
    title_slide = add_slide(presentation, 0)
    change_title(title_slide, title, textlayout=Textlayout(68, "Meiryo", PP_ALIGN.CENTER))
    title_slide.placeholders[1].text = subtitle


def change_title(slide, title, textlayout):
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].alignment = textlayout.textalign
    slide.shapes.title.text_frame.paragraphs[0].font.name = textlayout.fontfamily
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(textlayout.fontsize)

def add_text_slide(presentation, title, text, font_size, font_family, PP_ALIGN):
    slide = add_slide(presentation, 1)   # No1=title and content slide
    change_title(slide, title, textlayout=Textlayout(font_size, font_family, PP_ALIGN)) #テキストレイアウトの調整
    slide.placeholders[1].text = text    
    slide.placeholders[1].text_frame.paragraphs[0].font.size = Pt(26)
    slide.placeholders[1].text_frame.paragraphs[0].font.name = "Meiryo"

    return slide

def add_text_slide2(presentation, title, texts, font_size, font_family, PP_ALIGN):
    slide = add_slide(presentation, 1)   # No1=title and content slide
    change_title(slide, title, textlayout=Textlayout(font_size, font_family, PP_ALIGN)) #テキストレイアウトの調整
    for text in texts:
        new_paragraph = slide.shapes[1].text_frame.add_paragraph()  # cntent部分に段落を追加
        new_paragraph.text = text
        new_paragraph.alignment = PP_ALIGN #右揃えはRIGHT, 中央揃えはCENTER
        new_paragraph.font.bold = True  #太字の設定
        #new_paragraph.font.color.rgb = RGBColor(0,0,0) #文字の色設定
        new_paragraph.font.size = Pt(26)  #フォントサイズの設定
        new_paragraph.font.name = "Meiryo"
    return slide



def add_top_slide(presentation, title, text, font_size, font_family, PP_ALIGN):
    slide = add_slide(presentation, 1)   # No1=title and content slide
    change_title(slide, title, textlayout=Textlayout(font_size, font_family, PP_ALIGN)) #テキストレイアウトの調整
    slide.placeholders[1].text = text    
    slide.placeholders[1].text_frame.paragraphs[0].font.size = Pt(32)
    slide.placeholders[1].text_frame.paragraphs[0].font.name = "Meiryo"
    adjust_textbox(slide, Inches(10), Inches(5), Inches(3), Inches(2))
    return slide

def add_text_slide_with_background(presentation, title, text, background):
    slide = add_slide(presentation, 1)
    addImage(slide, background)
    change_title(slide, title, textlayout=Textlayout(48, "Dubai", PP_ALIGN.LEFT))
    slide.background
    slide.placeholders[1].text = text
    slide.placeholders[1].text_frame.paragraphs[0].font.size = Pt(32)
    return slide

def adjust_textbox(slide, left, top, width, height):
    textbox = slide.placeholders[1]
    textbox.left = left
    textbox.top = top
    textbox.width = width
    textbox.height = height
    

class Textlayout:
    """
    textレイアウトの設定の保持
    """
    def __init__(self, fontsize, fontfamily, textalign):
        self.fontsize = fontsize        # fontサイズ
        self.fontfamily = fontfamily    # fontの種類
        self.textalign = textalign      # textの位置